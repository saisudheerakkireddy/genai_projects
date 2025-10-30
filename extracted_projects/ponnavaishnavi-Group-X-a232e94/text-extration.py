import os

# Step 1: Upload
from google.colab import files
uploaded = files.upload()
filename = list(uploaded.keys())[0]

text = ""

# Step 2: Extract text based on file type
if filename.endswith(".txt"):
    with open(filename, "r", encoding="utf-8") as f:
        text = f.read()

elif filename.endswith(".pdf"):
    !pip install -q pymupdf
    import fitz
    with fitz.open(filename) as pdf:
        for page in pdf:
            text += page.get_text()

elif filename.endswith(".docx"):
    !pip install -q python-docx
    from docx import Document
    doc = Document(filename)
    text = "\n".join([para.text for para in doc.paragraphs])

elif filename.endswith(".csv"):
    import pandas as pd
    df = pd.read_csv(filename)
    text = df.to_string()

elif filename.endswith(".json"):
    import json, pprint
    with open(filename, "r", encoding="utf-8") as f:
        data = json.load(f)
    text = pprint.pformat(data)

elif filename.endswith(".pptx"):
    !pip install -q python-pptx
    from pptx import Presentation
    prs = Presentation(filename)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

elif filename.lower().endswith((".png", ".jpg", ".jpeg")):
    !pip install -q pytesseract pillow
    from PIL import Image
    import pytesseract
    img = Image.open(filename)
    text = pytesseract.image_to_string(img)

else:
    text = "❌ Unsupported file format."

# Step 3: View text
from IPython.display import display, HTML
if len(text.strip()) == 0:
    print("⚠ No text extracted.")
else:
    print("✅ Text extracted successfully!")
    display(HTML(f"<pre style='white-space: pre-wrap; font-size:14px;'>{text[:10000]}</pre>"))