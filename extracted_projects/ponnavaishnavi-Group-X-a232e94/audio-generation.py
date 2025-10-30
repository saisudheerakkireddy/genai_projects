# Full AI Explainer Audio Pipeline for Colab (Female Voice)

# Step 0: Install necessary packages
!pip install -q transformers torch gTTS pymupdf python-docx python-pptx pytesseract

# Step 1: Upload file
from google.colab import files
uploaded = files.upload()
filename = list(uploaded.keys())[0]

text = ""

# Step 2: Extract text based on file type
if filename.endswith(".txt"):
    with open(filename, "r", encoding="utf-8") as f:
        text = f.read()

elif filename.endswith(".pdf"):
    import fitz
    with fitz.open(filename) as pdf:
        for page in pdf:
            text += page.get_text()

elif filename.endswith(".docx"):
    from docx import Document
    doc = Document(filename)
    text = "\n".join([para.text for para in doc.paragraphs])

elif filename.endswith(".csv"):
    import pandas as pd
    df = pd.read_csv(filename)
    text = df.to_string()

elif filename.endswith(".json"):
    import json, pprint
    with open(filename, "r", encoding="utf-8") as f:
        data = json.load(f)
    text = pprint.pformat(data)

elif filename.endswith(".pptx"):
    from pptx import Presentation
    prs = Presentation(filename)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

elif filename.lower().endswith((".png", ".jpg", ".jpeg")):
    from PIL import Image
    import pytesseract
    img = Image.open(filename)
    text = pytesseract.image_to_string(img)

else:
    text = "❌ Unsupported file format."

# Step 3: Verify extracted text
from IPython.display import display, HTML

if len(text.strip()) == 0:
    print("⚠ No text extracted.")
    summary = "No content to summarize."
else:
    print("✅ Text extracted successfully!")
    display(HTML(f"<pre style='white-space: pre-wrap; font-size:14px;'>{text[:10000]}</pre>"))

    # Step 4: Summarize text with chunking
    from transformers import pipeline
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

    def chunk_text(text, chunk_size=500):
        words = text.split()
        for i in range(0, len(words), chunk_size):
            yield " ".join(words[i:i+chunk_size])

    summaries = []
    for chunk in chunk_text(text):
        summaries.append(summarizer(chunk, max_length=100, min_length=40, do_sample=False)[0]['summary_text'])

    summary = " ".join(summaries)

print("\n🧾 Summary:\n", summary)

# Step 5: Convert summary to audio (Female voice only using gTTS)
from gtts import gTTS
tts = gTTS(text=summary, lang='en')  # Female-like voice
audio_file = "summary_audio.mp3"
tts.save(audio_file)
print(f"\n🎤 Audio saved as {audio_file}")

# Step 6: Provide download link (Colab)
from google.colab import files
files.download(audio_file)