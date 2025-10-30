# ==========================================================
# 🎬 Professional AI Explainer Video Generator (Colab Ready)
# Corporate/HR Quality - Topic-Relevant Images
# Features: Professional animations, clear narration, business-ready
# ==========================================================

# STEP 1: Install all dependencies
print("📦 Installing required packages...")
!pip install -q transformers torch torchvision torchaudio diffusers accelerate safetensors
!pip install -q moviepy pillow gTTS pymupdf python-docx python-pptx pytesseract pandas requests
!pip install -q imageio==2.25.1 imageio-ffmpeg
!apt-get -qq install -y tesseract-ocr

print("✅ All packages installed!\n")

# STEP 2: Import all libraries
import os
import re
import textwrap
import torch
import requests
import json
import pandas as pd
import numpy as np
from PIL import Image, ImageDraw, ImageFont, ImageEnhance, ImageFilter
from gtts import gTTS
from diffusers import StableDiffusionPipeline
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips, TextClip, CompositeVideoClip
from IPython.display import display, HTML, Video as ipyVideo

print("✅ Libraries imported successfully!\n")

# ==========================================================
# STEP 3: Upload file or enter text manually
# ==========================================================
from google.colab import files

print("📁 Please upload your file (TXT, PDF, DOCX, PPTX, CSV, JSON, or Image)")
print("   Or press Cancel to enter text manually\n")

try:
    uploaded = files.upload()
    filename = list(uploaded.keys())[0]
    text = ""

    # Extract text based on file type
    if filename.endswith(".txt"):
        with open(filename, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

    elif filename.endswith(".pdf"):
        import fitz
        with fitz.open(filename) as pdf:
            for page in pdf:
                text += page.get_text()

    elif filename.endswith(".docx"):
        from docx import Document
        doc = Document(filename)
        text = "\n".join([para.text for para in doc.paragraphs])

    elif filename.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(filename)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif filename.endswith(".csv"):
        df = pd.read_csv(filename)
        text = df.to_string()

    elif filename.endswith(".json"):
        with open(filename, "r", encoding="utf-8") as f:
            data = json.load(f)
        text = json.dumps(data, indent=2)

    elif filename.lower().endswith((".png", ".jpg", ".jpeg")):
        import pytesseract
        img = Image.open(filename)
        text = pytesseract.image_to_string(img)

    else:
        print(f"⚠ File type '{filename}' not supported.")
        text = input("📝 Please enter your text manually:\n")

except:
    print("No file uploaded. Please enter text manually:")
    text = input("📝 Enter your text:\n")

# Validate extracted text
if len(text.strip()) == 0:
    raise ValueError("⚠ No text could be extracted. Please try again.")

print("\n✅ Text extracted successfully!")
print(f"📊 Length: {len(text)} characters\n")
display(HTML(f"<div style='background:#f0f0f0;padding:15px;border-radius:8px;max-height:300px;overflow:auto;'><pre style='white-space:pre-wrap;font-size:13px;'>{text[:2000]}{'...' if len(text) > 2000 else ''}</pre></div>"))

# ==========================================================
# STEP 4: Analyze topic and create key points
# ==========================================================
print("\n🧠 Analyzing content and extracting key points...")

from transformers import pipeline

# Load summarization model
summarizer = pipeline("summarization", model="facebook/bart-large-cnn", device=0 if torch.cuda.is_available() else -1)

def chunk_text(t, size=500):
    """Split text into chunks for processing"""
    words = t.split()
    for i in range(0, len(words), size):
        yield " ".join(words[i:i+size])

# Summarize text in chunks
summaries = []
chunks = list(chunk_text(text))
print(f"   Processing {len(chunks)} text chunks...")

for idx, chunk in enumerate(chunks, 1):
    if len(chunk.split()) < 30:
        continue
    try:
        result = summarizer(chunk, max_length=100, min_length=40, do_sample=False)[0]['summary_text']
        summaries.append(result)
        print(f"   ✓ Chunk {idx}/{len(chunks)} processed")
    except Exception as e:
        print(f"   ⚠ Skipped chunk {idx}: {str(e)[:50]}")

summary = " ".join(summaries)

# Break into key points
sentences = re.split(r'(?<=[.!?])\s+', summary)
sentences = [s.strip() for s in sentences if len(s.strip()) > 15][:10]  # Max 10 key points

if len(sentences) == 0:
    sentences = [text[:200]]  # Fallback

print(f"\n📝 Extracted {len(sentences)} key points:\n")
for i, s in enumerate(sentences, 1):
    print(f"   {i}. {s[:80]}{'...' if len(s) > 80 else ''}")

# ==========================================================
# STEP 5: Generate professional voiceover
# ==========================================================
print("\n🎤 Generating professional voiceover...")

narration_text = ". ".join(sentences) + "."
tts = gTTS(text=narration_text, lang='en', slow=False, tld='com')
tts.save("narration.mp3")

print("✅ Professional voiceover created")

# ==========================================================
# STEP 6: Setup professional fonts
# ==========================================================
print("\n📝 Setting up professional fonts...")

def download_font():
    """Download font with multiple fallback URLs"""
    font_urls = [
        "https://github.com/google/fonts/raw/main/apache/roboto/static/Roboto-Bold.ttf",
        "https://github.com/google/fonts/raw/main/ofl/opensans/static/OpenSans-Bold.ttf",
        "https://github.com/googlefonts/roboto/raw/main/src/hinted/Roboto-Bold.ttf"
    ]

    font_path = "CustomFont.ttf"

    for url in font_urls:
        try:
            print(f"   Trying: {url.split('/')[-1]}")
            response = requests.get(url, timeout=10)
            response.raise_for_status()

            with open(font_path, "wb") as f:
                f.write(response.content)

            # Test if font loads properly
            test_font = ImageFont.truetype(font_path, 20)
            print(f"✅ Font downloaded: {font_path}")
            return font_path

        except Exception as e:
            print(f"   ⚠ Failed: {str(e)[:50]}")
            continue

    # Ultimate fallback: use system default
    print("   Using system default font")
    return ImageFont.load_default()

font_path = download_font()

# ==========================================================
# STEP 7: Generate topic-relevant professional images
# ==========================================================
print("\n🎨 Generating professional, topic-relevant images...")
print("   (This may take 5-10 minutes)\n")

device = "cuda" if torch.cuda.is_available() else "cpu"
print(f"   Using device: {device}")

# Load Stable Diffusion pipeline
pipe = StableDiffusionPipeline.from_pretrained(
    "runwayml/stable-diffusion-v1-5",
    torch_dtype=torch.float16 if device == "cuda" else torch.float32,
    safety_checker=None
).to(device)

# Optimize memory
if device == "cuda":
    pipe.enable_attention_slicing()
    pipe.enable_vae_slicing()

# Create directories
os.makedirs("frames", exist_ok=True)
os.makedirs("animated", exist_ok=True)

# Extract topic keywords for better prompts
def extract_topic_context(text_content):
    """Extract main topic and context from content"""
    keywords = text_content.lower()

    topic_hints = {
        'business': ['business', 'company', 'corporate', 'management', 'strategy', 'market'],
        'hr': ['employee', 'hr', 'human resources', 'recruitment', 'talent', 'team', 'workplace'],
        'technology': ['technology', 'software', 'digital', 'ai', 'computer', 'data'],
        'finance': ['finance', 'financial', 'money', 'investment', 'accounting', 'budget'],
        'healthcare': ['health', 'medical', 'patient', 'clinical', 'hospital', 'care'],
        'education': ['education', 'learning', 'training', 'student', 'teaching', 'course'],
        'sales': ['sales', 'customer', 'revenue', 'client', 'marketing', 'product'],
        'operations': ['operations', 'process', 'logistics', 'supply', 'production', 'workflow']
    }

    detected_topics = []
    for topic, terms in topic_hints.items():
        if any(term in keywords for term in terms):
            detected_topics.append(topic)

    return detected_topics[0] if detected_topics else 'professional business'

main_topic = extract_topic_context(text)
print(f"   Detected topic: {main_topic.upper()}\n")

image_paths = []

for i, sentence in enumerate(sentences):
    print(f"   🖼  Generating scene {i+1}/{len(sentences)}...")

    # Create professional, topic-relevant prompt
    prompt = (
        f"professional {main_topic} concept, {sentence[:100]}, "
        f"clean modern office setting, business environment, "
        f"photorealistic, high quality, corporate style, "
        f"professional photography, bright natural lighting, "
        f"sharp focus, detailed, 4k quality"
    )

    negative_prompt = (
        "cartoon, anime, illustration, drawing, painting, sketch, "
        "low quality, blurry, dark, amateur, text, watermark, "
        "unprofessional, messy, cluttered"
    )

    # Generate image
    try:
        img = pipe(
            prompt,
            negative_prompt=negative_prompt,
            num_inference_steps=40,
            guidance_scale=7.5,
            height=1024,
            width=1024
        ).images[0]

        # Professional post-processing
        # Slight blur for softer look
        img = img.filter(ImageFilter.GaussianBlur(radius=0.5))

        # Enhance for corporate look
        enhancer = ImageEnhance.Contrast(img)
        img = enhancer.enhance(1.1)
        enhancer = ImageEnhance.Brightness(img)
        img = enhancer.enhance(1.05)
        enhancer = ImageEnhance.Sharpness(img)
        img = enhancer.enhance(1.2)

        # Resize to HD 16:9
        img = img.resize((1920, 1080), Image.Resampling.LANCZOS)

        path = f"frames/scene_{i:03d}.png"
        img.save(path, quality=95)
        image_paths.append(path)
        print(f"      ✓ Professional scene {i+1} created\n")

    except Exception as e:
        print(f"      ⚠ Error: {str(e)[:50]}")
        # Create professional gradient background fallback
        fallback = Image.new('RGB', (1920, 1080), color=(25, 35, 55))
        draw = ImageDraw.Draw(fallback)
        # Add subtle gradient effect
        for y in range(1080):
            color_val = int(25 + (y / 1080) * 30)
            draw.line([(0, y), (1920, y)], fill=(color_val, color_val+10, color_val+30))

        path = f"frames/scene_{i:03d}.png"
        fallback.save(path)
        image_paths.append(path)

print(f"✅ Generated {len(image_paths)} professional scenes\n")

# ==========================================================
# STEP 8: Add professional text overlays
# ==========================================================
print("📝 Adding professional text overlays...")

processed_paths = []

for i, path in enumerate(image_paths):
    img = Image.open(path).convert("RGB")

    # Apply professional vignette effect
    vignette = Image.new('RGBA', img.size, (0, 0, 0, 0))
    vignette_draw = ImageDraw.Draw(vignette)

    width, height = img.size
    for y in range(height):
        for x in range(width):
            # Calculate distance from center
            dx = (x - width/2) / (width/2)
            dy = (y - height/2) / (height/2)
            distance = (dx*2 + dy*2) ** 0.5

            if distance > 0.7:
                alpha = int((distance - 0.7) * 200)
                alpha = min(alpha, 100)
                vignette_draw.point((x, y), fill=(0, 0, 0, alpha))

    img = img.convert('RGBA')
    img = Image.alpha_composite(img, vignette)

    # Create professional text overlay
    overlay = Image.new('RGBA', img.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)

    # Prepare text with proper font handling
    try:
        if isinstance(font_path, str):
            font = ImageFont.truetype(font_path, 52)
        else:
            font = font_path
    except:
        font = ImageFont.load_default()

    wrapped_text = textwrap.fill(sentences[i], width=55)

    # Calculate text dimensions
    bbox = draw.multiline_textbbox((0, 0), wrapped_text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]

    # Draw professional lower-third style background
    padding = 50
    bg_height = text_height + padding * 2
    bg_y = img.height - bg_height - 40

    # Gradient background for text
    for offset in range(bg_height):
        alpha = int(160 + (offset / bg_height) * 40)
        draw.rectangle(
            [(0, bg_y + offset), (img.width, bg_y + offset + 1)],
            fill=(10, 20, 40, min(alpha, 200))
        )

    # Add accent line (corporate blue)
    draw.rectangle(
        [(0, bg_y), (img.width, bg_y + 4)],
        fill=(0, 120, 215, 255)
    )

    # Draw text
    text_x = (img.width - text_width) / 2
    text_y = bg_y + padding
    draw.multiline_text(
        (text_x, text_y),
        wrapped_text,
        font=font,
        fill=(255, 255, 255, 255),
        align="center"
    )

    # Composite overlay
    img = Image.alpha_composite(img, overlay)

    # Save processed image
    processed_path = f"animated/scene_{i:03d}.png"
    img.convert('RGB').save(processed_path, quality=95)
    processed_paths.append(processed_path)
    print(f"   ✓ Scene {i+1}/{len(image_paths)} styled")

print("\n✅ All professional overlays added\n")

# ==========================================================
# STEP 9: Create smooth professional animations
# ==========================================================
print("🎬 Creating professional video with smooth animations...")

audio = AudioFileClip("narration.mp3")
dur_per_scene = audio.duration / len(processed_paths)
print(f"   Duration per scene: {dur_per_scene:.2f} seconds")

clips = []

for i, path in enumerate(processed_paths):
    print(f"   🎞  Animating scene {i+1}/{len(processed_paths)}...")

    # Create base clip
    clip = ImageClip(path, duration=dur_per_scene)

    # Subtle zoom (professional, not overwhelming)
    def make_zoom(t):
        return 1.0 + 0.05 * (t / dur_per_scene)

    clip = clip.resize(make_zoom)

    # Smooth professional transitions
    clip = clip.crossfadein(0.8).crossfadeout(0.8)

    clips.append(clip)

print("\n✅ All animations created\n")

# ==========================================================
# STEP 10: Create professional intro and outro
# ==========================================================
print("📺 Creating professional intro and outro...")

def create_professional_card(title, subtitle, duration=3.5, is_intro=True):
    """Create professional title card with gradient"""
    # Corporate gradient background
    img = Image.new('RGB', (1920, 1080))
    draw = ImageDraw.Draw(img)

    # Professional gradient (dark blue to deep navy)
    for y in range(1080):
        r = int(15 + (y / 1080) * 20)
        g = int(25 + (y / 1080) * 30)
        b = int(45 + (y / 1080) * 40)
        draw.line([(0, y), (1920, y)], fill=(r, g, b))

    # Add geometric accent
    overlay = Image.new('RGBA', img.size, (0, 0, 0, 0))
    overlay_draw = ImageDraw.Draw(overlay)

    # Diagonal accent stripe
    overlay_draw.rectangle(
        [(0, 450), (1920, 460)],
        fill=(0, 120, 215, 180)
    )

    img = img.convert('RGBA')
    img = Image.alpha_composite(img, overlay)
    draw = ImageDraw.Draw(img)

    # Load fonts with fallback
    try:
        if isinstance(font_path, str):
            title_font = ImageFont.truetype(font_path, 95)
            subtitle_font = ImageFont.truetype(font_path, 48)
        else:
            title_font = ImageFont.load_default()
            subtitle_font = ImageFont.load_default()
    except:
        title_font = ImageFont.load_default()
        subtitle_font = ImageFont.load_default()

    # Draw title
    title_bbox = draw.textbbox((0, 0), title, font=title_font)
    title_w = title_bbox[2] - title_bbox[0]
    draw.text(
        ((1920 - title_w) / 2, 380),
        title,
        font=title_font,
        fill=(255, 255, 255, 255)
    )

    # Draw subtitle
    sub_bbox = draw.textbbox((0, 0), subtitle, font=subtitle_font)
    sub_w = sub_bbox[2] - sub_bbox[0]
    draw.text(
        ((1920 - sub_w) / 2, 550),
        subtitle,
        font=subtitle_font,
        fill=(180, 200, 230, 255)
    )

    # Save and create clip
    temp_path = f"animated/{'intro' if is_intro else 'outro'}.png"
    img.convert('RGB').save(temp_path, quality=95)

    clip = ImageClip(temp_path, duration=duration).fadein(1.2).fadeout(1.2)
    return clip

intro = create_professional_card("Professional Explainer", "Key Insights Ahead", 3.5, True)
outro = create_professional_card("Thank You", "For Your Attention", 3.5, False)

print("✅ Professional intro and outro created\n")

# ==========================================================
# STEP 11: Assemble final professional video
# ==========================================================
print("🎥 Assembling final professional video...")

# Combine all clips
all_clips = [intro] + clips + [outro]
final_video = concatenate_videoclips(all_clips, method="compose")

# Add audio (start after intro)
final_video = final_video.set_audio(audio.set_start(3.5))

# Export with high quality settings
print("   Rendering HD video (this may take several minutes)...\n")

final_video.write_videofile(
    "Professional_Explainer_Video.mp4",
    fps=30,
    codec='libx264',
    audio_codec='aac',
    bitrate='8000k',
    preset='medium',
    threads=4
)

print("\n" + "="*70)
print("🎉 SUCCESS! Your professional explainer video is ready!")
print("="*70)
print("✅ File: Professional_Explainer_Video.mp4")
print(f"✅ Duration: {final_video.duration:.1f} seconds")
print(f"✅ Scenes: {len(sentences)}")
print("✅ Quality: HD 1920x1080, 30 FPS")
print("✅ Style: Professional, topic-relevant, corporate-ready")
print("✅ Features: Smooth animations, clear narration, lower-thirds")
print("="*70 + "\n")

# Display video in Colab
print("📺 Playing video preview:")
ipyVideo("Professional_Explainer_Video.mp4", embed=True, width=800, height=450)

# Download button
print("\n💾 Download your professional video:")
files.download("Professional_Explainer_Video.mp4")